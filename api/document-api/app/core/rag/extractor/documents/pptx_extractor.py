"""Abstract interface for PowerPoint document loader implementations."""
import os
import subprocess
import time
from typing import Optional, List
from pptx import Presentation
from app.core.rag.extractor.extractor_base import BaseExtractor
from app.core.rag.models.document import Document


class PowerPointExtractor(BaseExtractor):
    """Load PowerPoint files (.ppt and .pptx).

    Args:
        file_path: Path to the file to load.
        include_slide_numbers: Whether to include slide numbers in extracted content.
        include_notes: Whether to include presenter notes in extracted content.
        conversion_tool: Tool to use for PPT conversion ('libreoffice' or 'none').
            If 'none', .ppt files will raise an error.
    """

    def __init__(
            self,
            file_path: str,
            include_slide_numbers: bool = True,
            include_notes: bool = True,
            conversion_tool: str = 'libreoffice',
    ):
        """Initialize with file path."""
        self._file_path = os.path.abspath(file_path)
        self._include_slide_numbers = include_slide_numbers
        self._include_notes = include_notes
        self._conversion_tool = conversion_tool
        self._temp_pptx_path = None

    def _convert_ppt_to_pptx_libreoffice(self) -> str:
        """Convert .ppt to .pptx format using LibreOffice."""
        try:
            # Create full paths
            base_dir = os.path.dirname(self._file_path)
            base_name = os.path.splitext(os.path.basename(self._file_path))[0]
            temp_path = os.path.join(base_dir, f"{base_name}.pptx")
            
            # Remove existing temp file if it exists
            if os.path.exists(temp_path):
                os.remove(temp_path)
            
            # Convert using LibreOffice
            process = subprocess.run([
                'libreoffice',  # Try libreoffice first
                '--headless',
                '--convert-to',
                'pptx',
                '--outdir',
                base_dir,
                self._file_path
            ], capture_output=True, text=True)
            
            if process.returncode != 0:
                # Try soffice if libreoffice command fails
                process = subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to',
                    'pptx',
                    '--outdir',
                    base_dir,
                    self._file_path
                ], capture_output=True, text=True)
            
            if process.returncode != 0:
                raise RuntimeError(f"LibreOffice conversion failed: {process.stderr}")
            
            # Wait for the file to be created (max 10 seconds)
            max_wait = 10
            wait_time = 0
            while not os.path.exists(temp_path) and wait_time < max_wait:
                time.sleep(0.5)
                wait_time += 0.5
            
            if not os.path.exists(temp_path):
                raise RuntimeError(f"Converted file not found at {temp_path}")
            
            # Verify the file is readable
            if not os.path.getsize(temp_path):
                raise RuntimeError("Converted file is empty")
                
            return temp_path
            
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"Error executing LibreOffice: {str(e)}")
        except FileNotFoundError:
            raise RuntimeError("LibreOffice not found. Please install LibreOffice or use a .pptx file.")

    def _extract_text_from_shape(self, shape) -> str:
        """Extract text from a shape, handling different shape types."""
        text = ""
        
        try:
            if hasattr(shape, "text"):
                text = shape.text
            elif hasattr(shape, "table"):
                # Extract text from tables
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        text += " | ".join(row_texts) + "\n"
            elif hasattr(shape, "chart"):
                # Extract title and labels from charts
                chart = shape.chart
                if hasattr(chart, "chart_title") and chart.chart_title and chart.chart_title.text_frame.text:
                    text += f"Chart: {chart.chart_title.text_frame.text}\n"
        except Exception as e:
            # Log error but continue processing other shapes
            print(f"Error extracting text from shape: {str(e)}")
            
        return text.strip()

    def _extract_pptx_content(self, file_path: str) -> List[Document]:
        """Extract content from .pptx file."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found at {file_path}")
            
        if not os.path.getsize(file_path):
            raise RuntimeError(f"PPTX file is empty: {file_path}")
            
        try:
            presentation = Presentation(file_path)
        except Exception as e:
            raise RuntimeError(f"Error opening PPTX file {file_path}: {str(e)}")

        documents = []
        
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_content = []
            
            # Add slide number if enabled
            if self._include_slide_numbers:
                slide_content.append(f"Slide {slide_number}")

            # Extract text from shapes
            for shape in slide.shapes:
                text = self._extract_text_from_shape(shape)
                if text:
                    slide_content.append(text)

            # Extract notes if enabled
            if self._include_notes and hasattr(slide, "notes_slide") and slide.notes_slide:
                for shape in slide.notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(f"Notes: {shape.text.strip()}")

            if slide_content:
                content = "\n\n".join(slide_content)
                documents.append(Document(page_content=content))

        return documents

    def extract(self) -> List[Document]:
        """Load from file path."""
        if not os.path.exists(self._file_path):
            raise FileNotFoundError(f"File not found: {self._file_path}")

        try:
            # Check if file is PPT or PPTX
            if self._file_path.lower().endswith('.ppt'):
                if self._conversion_tool == 'libreoffice':
                    self._temp_pptx_path = self._convert_ppt_to_pptx_libreoffice()
                    return self._extract_pptx_content(self._temp_pptx_path)
                else:
                    raise ValueError("Legacy .ppt files are only supported with LibreOffice conversion. "
                                  "Please convert the file to .pptx format or enable LibreOffice conversion.")
            elif self._file_path.lower().endswith('.pptx'):
                return self._extract_pptx_content(self._file_path)
            else:
                raise ValueError("Unsupported file format. Only .ppt and .pptx files are supported.")

        finally:
            # Cleanup temporary file if it was created
            if self._temp_pptx_path and os.path.exists(self._temp_pptx_path):
                try:
                    os.remove(self._temp_pptx_path)
                except Exception:
                    pass  # Ignore cleanup errors